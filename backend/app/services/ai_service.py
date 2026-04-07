"""AI Service â€” Groq-powered AI features for students.

Uses Groq API with Llama / Mixtral models for:
  - Video summary, key points, follow-up questions
  - Document important-points extraction
  - Course-specific and general chatbot
  - Speech-to-text via Whisper
"""

import os
import re
import json
import requests
from pathlib import Path
from typing import Optional

from deep_translator import GoogleTranslator

GROQ_API_KEY = os.getenv("GROQ_API_KEY")
GROQ_CHAT_URL = "https://api.groq.com/openai/v1/chat/completions"
GROQ_TRANSCRIPTION_URL = "https://api.groq.com/openai/v1/audio/transcriptions"

# Fast + free model on Groq
CHAT_MODEL = "llama-3.3-70b-versatile"
FAST_MODEL = "llama-3.1-8b-instant"
WHISPER_MODEL = "whisper-large-v3"

# Maps frontend language code â†’ (display name, Google Translate ISO code)
SUPPORTED_TRANSLATION_LANGS: dict[str, str] = {
    "en": "english",
    "hi": "hindi",
    "te": "telugu",
    "ta": "tamil",
    "kn": "kannada",
    "ml": "malayalam",
    "mr": "marathi",
    "bn": "bengali",
    "gu": "gujarati",
    "pa": "punjabi",
    "ur": "urdu",
}

# Maps full name â†’ ISO code for reverse lookup (used by translate_text)
_LANG_NAME_TO_CODE: dict[str, str] = {v: k for k, v in SUPPORTED_TRANSLATION_LANGS.items()}


def _headers():
    return {
        "Authorization": f"Bearer {GROQ_API_KEY}",
        "Content-Type": "application/json",
    }


def _chat(messages: list[dict], model: str = CHAT_MODEL, temperature: float = 0.4, max_tokens: int = 4096) -> str:
    """Send a chat completion request to Groq and return the assistant text."""
    payload = {
        "model": model,
        "messages": messages,
        "temperature": temperature,
        "max_tokens": max_tokens,
    }
    resp = requests.post(GROQ_CHAT_URL, headers=_headers(), json=payload, timeout=60)
    resp.raise_for_status()
    data = resp.json()
    return data["choices"][0]["message"]["content"]


def _chunk_text(text: str, chunk_size: int = 4200) -> list[str]:
    """Split text into translation-safe chunks (Google web translator has payload limits)."""
    value = (text or "").strip()
    if not value:
        return []
    if len(value) <= chunk_size:
        return [value]

    chunks: list[str] = []
    current: list[str] = []
    current_len = 0

    for sentence in re.split(r"(?<=[.!?])\s+", value):
        sentence = sentence.strip()
        if not sentence:
            continue

        if len(sentence) > chunk_size:
            if current:
                chunks.append(" ".join(current).strip())
                current = []
                current_len = 0
            for i in range(0, len(sentence), chunk_size):
                part = sentence[i:i + chunk_size].strip()
                if part:
                    chunks.append(part)
            continue

        projected = current_len + len(sentence) + (1 if current else 0)
        if projected > chunk_size and current:
            chunks.append(" ".join(current).strip())
            current = [sentence]
            current_len = len(sentence)
        else:
            current.append(sentence)
            current_len = projected

    if current:
        chunks.append(" ".join(current).strip())
    return chunks


def translate_text(text: str, target_language: str, source_language: str = "auto") -> str:
    """Translate text with deep-translator Google backend."""
    if not text or not text.strip():
        return text

    raw = (target_language or "").strip().lower()

    # Normalise: accept both ISO code ("hi") and full name ("hindi")
    if raw in ("", "en", "english"):
        return text

    # Resolve to ISO code â€” GoogleTranslator requires codes, not full names
    if raw in SUPPORTED_TRANSLATION_LANGS:
        iso_code = raw                        # already a code e.g. "hi"
    elif raw in _LANG_NAME_TO_CODE:
        iso_code = _LANG_NAME_TO_CODE[raw]    # full name â†’ code e.g. "hindi" â†’ "hi"
    else:
        raise ValueError(f"Unsupported target language: {target_language!r}")

    translator = GoogleTranslator(source=source_language, target=iso_code)
    chunks = _chunk_text(text)
    if not chunks:
        return ""

    translated_chunks: list[str] = []
    for chunk in chunks:
        result = translator.translate(chunk)
        translated_chunks.append((result or "").strip())
    return "\n".join(part for part in translated_chunks if part)


def _translate_list(items: list[str], target_language: str) -> list[str]:
    """Translate a list of strings concurrently using threads."""
    from concurrent.futures import ThreadPoolExecutor, as_completed

    if not items:
        return []

    results: list[str] = [""] * len(items)

    def _do(idx: int, text: str) -> tuple[int, str]:
        return idx, translate_text(text, target_language)

    with ThreadPoolExecutor(max_workers=6) as pool:
        futures = {pool.submit(_do, i, t): i for i, t in enumerate(items)}
        for fut in as_completed(futures):
            i, translated = fut.result()
            results[i] = translated

    return results


def translate_flashcards(cards: list, target_language: str) -> list:
    """Translate a list of flashcard front/back pairs."""
    translated = []
    for card in cards:
        front = translate_text(str(card.get("front", "")), target_language)
        back = translate_text(str(card.get("back", "")), target_language)
        translated.append({"front": front, "back": back})
    return translated


def translate_analysis(analysis: dict, analysis_type: str, target_language: str) -> dict:
    """Translate video/document analysis payload while preserving schema."""
    if not isinstance(analysis, dict):
        return {}

    translated = dict(analysis)
    translated["summary"] = translate_text(str(analysis.get("summary", "")), target_language)

    if analysis_type == "video":
        translated["key_points"] = _translate_list(
            [str(p) for p in analysis.get("key_points", []) if str(p).strip()],
            target_language,
        )
        translated["follow_up_questions"] = _translate_list(
            [str(q) for q in analysis.get("follow_up_questions", []) if str(q).strip()],
            target_language,
        )
    else:
        translated["important_points"] = _translate_list(
            [str(p) for p in analysis.get("important_points", []) if str(p).strip()],
            target_language,
        )
        # key_definitions: translate terms and definitions concurrently together
        defs = [d for d in analysis.get("key_definitions", []) if isinstance(d, dict)]
        terms = _translate_list([str(d.get("term", "")) for d in defs], target_language)
        definitions = _translate_list([str(d.get("definition", "")) for d in defs], target_language)
        translated["key_definitions"] = [
            {"term": terms[i], "definition": definitions[i]}
            for i in range(len(defs))
        ]
        translated["study_tips"] = _translate_list(
            [str(t) for t in analysis.get("study_tips", []) if str(t).strip()],
            target_language,
        )

    translated["translated_to"] = target_language
    return translated


# ---------------------------------------------------------------------------
# Speech-to-text  (video / audio transcription via Whisper on Groq)
# ---------------------------------------------------------------------------

def transcribe_audio(file_path: str, language: str = "en") -> str:
    """Transcribe an audio/video file using Groq Whisper API."""
    headers = {"Authorization": f"Bearer {GROQ_API_KEY}"}

    with open(file_path, "rb") as f:
        files = {"file": (Path(file_path).name, f, "audio/mpeg")}
        data = {
            "model": WHISPER_MODEL,
            "language": language,
            "response_format": "verbose_json",
        }
        resp = requests.post(
            GROQ_TRANSCRIPTION_URL, headers=headers, files=files, data=data, timeout=120
        )
    resp.raise_for_status()
    result = resp.json()
    return result.get("text", "")


# ---------------------------------------------------------------------------
# Video AI â€” summary, key points, follow-up questions
# ---------------------------------------------------------------------------

def analyze_video_content(transcript: str, video_title: str = "") -> dict:
    """Given a transcript, produce summary + key points + follow-up questions."""
    system = (
        "You are an expert educational AI assistant. "
        "Given a video transcript, produce a structured analysis. "
        "Return ONLY valid JSON with these exact keys:\n"
        '  "summary": a concise 3-5 sentence summary,\n'
        '  "key_points": an array of 5-8 important points (strings),\n'
        '  "follow_up_questions": an array of 3-5 thought-provoking follow-up questions for students.\n'
        "Do NOT include markdown code fences or any other text outside the JSON."
    )
    user_msg = f"Video title: {video_title}\n\nTranscript:\n{transcript[:12000]}"

    raw = _chat([
        {"role": "system", "content": system},
        {"role": "user", "content": user_msg},
    ], temperature=0.3)

    return _parse_json(raw, fallback={
        "summary": raw,
        "key_points": [],
        "follow_up_questions": [],
    })


def summarize_video_from_url(video_url: str, video_title: str = "") -> dict:
    """For videos we cannot transcribe (e.g. YouTube links), generate analysis from title + URL context."""
    system = (
        "You are an expert educational AI assistant. "
        "A student is watching a video and wants help understanding it. "
        "Based on the video title and any available context, provide a helpful "
        "educational analysis. Return ONLY valid JSON with these exact keys:\n"
        '  "summary": a helpful overview of what this video likely covers (3-5 sentences),\n'
        '  "key_points": an array of 5-8 key concepts students should look for,\n'
        '  "follow_up_questions": an array of 3-5 study questions related to the topic.\n'
        "Do NOT include markdown code fences or any other text outside the JSON."
    )
    user_msg = f"Video title: {video_title}\nVideo URL: {video_url}"

    raw = _chat([
        {"role": "system", "content": system},
        {"role": "user", "content": user_msg},
    ], temperature=0.4)

    return _parse_json(raw, fallback={
        "summary": raw,
        "key_points": [],
        "follow_up_questions": [],
    })


# ---------------------------------------------------------------------------
# Document AI â€” important points extraction
# ---------------------------------------------------------------------------

def analyze_document(text_content: str, doc_title: str = "") -> dict:
    """Extract important points and a summary from document text."""
    system = (
        "You are an expert educational AI assistant. "
        "Given the text content of a document, produce a structured analysis. "
        "Return ONLY valid JSON with these exact keys:\n"
        '  "summary": a concise 3-5 sentence summary of the document,\n'
        '  "important_points": an array of 6-10 important points (strings),\n'
        '  "key_definitions": an array of objects with "term" and "definition" keys for important terms,\n'
        '  "study_tips": an array of 3-5 practical study tips related to this content.\n'
        "Do NOT include markdown code fences or any other text outside the JSON."
    )
    user_msg = f"Document title: {doc_title}\n\nContent:\n{text_content[:12000]}"

    raw = _chat([
        {"role": "system", "content": system},
        {"role": "user", "content": user_msg},
    ], temperature=0.3)

    return _parse_json(raw, fallback={
        "summary": raw,
        "important_points": [],
        "key_definitions": [],
        "study_tips": [],
    })


def analyze_document_from_url(doc_url: str, doc_title: str = "") -> dict:
    """For documents we can't read directly, provide analysis from title/URL context."""
    system = (
        "You are an expert educational AI assistant. "
        "A student has a study document and wants help understanding it. "
        "Based on the document title and context, provide a helpful educational analysis. "
        "Return ONLY valid JSON with these exact keys:\n"
        '  "summary": a helpful overview of what this document likely covers (3-5 sentences),\n'
        '  "important_points": an array of 6-10 key concepts students should focus on,\n'
        '  "key_definitions": an array of objects with "term" and "definition" keys for likely important terms,\n'
        '  "study_tips": an array of 3-5 practical study tips related to this content.\n'
        "Do NOT include markdown code fences or any other text outside the JSON."
    )
    user_msg = f"Document title: {doc_title}\nDocument URL: {doc_url}"

    raw = _chat([
        {"role": "system", "content": system},
        {"role": "user", "content": user_msg},
    ], temperature=0.4)

    return _parse_json(raw, fallback={
        "summary": raw,
        "important_points": [],
        "key_definitions": [],
        "study_tips": [],
    })


# ---------------------------------------------------------------------------
# Chatbot â€” general + course-specific
# ---------------------------------------------------------------------------

def chat_with_ai(
    messages: list[dict],
    course_context: Optional[str] = None,
    system_extra: str = "",
) -> str:
    """General-purpose chatbot. Optionally inject course context."""
    system_parts = [
        "You are a helpful, friendly educational AI assistant for a college student learning platform.",
        "You help students understand their course materials, answer questions, and provide study guidance.",
        "Be concise but thorough. Use examples when helpful. Format your responses with markdown for readability.",
    ]
    if course_context:
        system_parts.append(
            f"\n\nCourse Context (use this to give course-specific answers):\n{course_context}"
        )
    if system_extra:
        system_parts.append(system_extra)

    system_msg = {"role": "system", "content": " ".join(system_parts)}
    all_messages = [system_msg] + messages[-20:]  # keep last 20 messages for context

    return _chat(all_messages, model=CHAT_MODEL, temperature=0.5, max_tokens=2048)


def build_course_context(course_data: dict) -> str:
    """Build a text context string from course data for the chatbot."""
    parts = [
        f"Course: {course_data.get('name', 'Unknown')}",
        f"Description: {course_data.get('description', 'N/A')}",
        f"Teacher: {course_data.get('teacher', {}).get('name', 'Unknown')}",
    ]

    resources = course_data.get("resources", [])
    if resources:
        parts.append(f"\nCourse Resources ({len(resources)} total):")
        for i, r in enumerate(resources, 1):
            line = f"  {i}. {r.get('title', 'Untitled')} ({r.get('type', 'unknown')})"
            if r.get("notes"):
                line += f" â€” Notes: {r['notes']}"
            parts.append(line)

    return "\n".join(parts)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _parse_json(raw: str, fallback: dict) -> dict:
    """Try to parse JSON from the model output, with fallbacks."""
    # Strip markdown code fences if present
    cleaned = re.sub(r"```json\s*", "", raw)
    cleaned = re.sub(r"```\s*", "", cleaned)
    cleaned = cleaned.strip()

    try:
        return json.loads(cleaned)
    except json.JSONDecodeError:
        # Try to find JSON object in the response
        match = re.search(r"\{.*\}", cleaned, re.DOTALL)
        if match:
            try:
                return json.loads(match.group())
            except json.JSONDecodeError:
                pass
        return fallback


def read_text_file(file_path: str) -> str:
    """Read text content from common file formats."""
    path = Path(file_path)
    ext = path.suffix.lower()

    # Plain text files
    if ext in (".txt", ".md", ".csv", ".json", ".xml", ".yaml", ".yml", ".log", ".rtf"):
        try:
            return path.read_text(encoding="utf-8", errors="ignore")
        except Exception:
            return ""

    # PDF extraction
    if ext == ".pdf":
        try:
            import PyPDF2
            text_parts = []
            with open(file_path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages[:50]:  # limit pages
                    text_parts.append(page.extract_text() or "")
            return "\n".join(text_parts)
        except ImportError:
            return "[PDF reading requires PyPDF2 package]"
        except Exception:
            return ""

    # DOCX extraction
    if ext in (".docx",):
        try:
            import docx
            doc = docx.Document(file_path)
            return "\n".join(p.text for p in doc.paragraphs)
        except ImportError:
            return "[DOCX reading requires python-docx package]"
        except Exception:
            return ""

    # PPTX extraction
    if ext in (".pptx",):
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            text_parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_parts.append(shape.text)
            return "\n".join(text_parts)
        except ImportError:
            return "[PPTX reading requires python-pptx package]"
        except Exception:
            return ""

    return ""
